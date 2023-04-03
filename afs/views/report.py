from io import BytesIO
from types import SimpleNamespace
import logging
import json
import re
import uuid
from django.conf import settings
from django.http import HttpResponse, HttpResponseServerError
from django.utils.translation import ugettext as _
from django.views.generic import View
from django.utils.translation import get_language
from arches.app.utils.response import JSONResponse
from arches.app.models.models import PublishedGraph, Node
from arches.app.models.resource import Resource
from arches.app.datatypes.datatypes import DataTypeFactory
from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


class ReportView(View):
    def get(self, request):
        return JSONResponse(settings.AFS_CUSTOM_REPORTS)

    def post(self, request):
        json_data = json.loads(request.body)
        title_key = "b5b38bde-b2f2-11e9-aff6-a4d18cec433a"
        template_id = json_data["templateid"] if "templateid" in json_data else None
        project_id = json_data["projectid"] if "projectid" in json_data else None
        template = settings.AFS_CUSTOM_REPORTS[template_id] if template_id in settings.AFS_CUSTOM_REPORTS else None
        project = Resource.objects.filter(pk=project_id)[0]
        published_graph = PublishedGraph.objects.filter(publication_id=str(project.graph_publication_id), language=get_language()).first()

        datatype_factory = DataTypeFactory()
        serialized_graph = published_graph.serialized_graph

        node_dict = {}

        for index, i in enumerate(serialized_graph["nodes"]):
            if serialized_graph["nodes"][index]["alias"]:
                serialized_graph["nodes"][index]["pk"] = uuid.UUID(serialized_graph["nodes"][index]["nodeid"])
                node_dict[serialized_graph["nodes"][index]["alias"].lower()] = SimpleNamespace(**serialized_graph["nodes"][index])

        project.load_tiles()
        tiles = project.tiles

        bytestream = BytesIO()

        if template.endswith("docx"):
            doc = Document(template)

            title = str(project.name)

            for p in doc.paragraphs:
                arches_tag_pattern = re.compile(r"(\<arches\:\s?([A-Za-z_0-9]+)\>)")  # should match <arches: node_alias>
                for run in p.runs:
                    for (tag, alias) in re.findall(arches_tag_pattern, run.text):
                        tile = next((d for d in project.tiles if str(d.nodegroup_id) == node_dict[alias].nodegroup_id), None)
                        if tile is not None:
                            node = node_dict[alias]
                            display_value = datatype_factory.get_instance(node.datatype).get_display_value(tile, node)
                            run.text = run.text.replace(tag, display_value)

            doc.save(bytestream)
            mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        elif template.endswith("pptx"):
            pptx = Presentation(template)

            for tile in tiles:
                for key in tile.data:
                    if key == title_key:
                        title = tile.data[key] + ".pptx"
                    for s in pptx.slides:
                        for shape in s.shapes:
                            if key in shape.text:
                                shape.text = shape.text.replace("<" + key + ">", tile.data[key])
            pptx.save(bytestream)
            mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        bytestream.seek(0)

        if template is None:
            return HttpResponseServerError("Could not find requested template")

        response = HttpResponse(content=bytestream.read())
        response["Content-Type"] = mime
        response["Content-Disposition"] = "attachment; filename={}".format(title)
        return response


# class ArchesTemplateEngine():

#     def __init__(self, language=None):
#         if not language:
#             self.language = get_language()
#         else:
#             self.language = language

#         published_graphs = PublishedGraph.objects.filter(language=get_language())

#         self.serialized_graphs = {}
#         self.patterns = [
#             {
#                 "pattern": re.compile(r'(\<arches\:\s?([A-Za-z_0-9]+)\>)'),
#                 "replacement_method": "single"
#             }
#         ]

#         for graph in published_graphs:
#             serialized_graph = graph.serialized_graph

#             for index, i in enumerate(serialized_graph['nodes']):
#                 if serialized_graph['nodes'][index]["alias"]:
#                     serialized_graph['nodes'][index]["pk"] = uuid.UUID(serialized_graph['nodes'][index]["nodeid"])
#                     serialized_graph['node_namespace'] = SimpleNamespace(**serialized_graph['nodes'][index])

#             self.serialized_graphs[graph.graphid] = serialized_graph

#     def document_replace(self, document, resources):
#         pass

#     class DocxTemplateEngine(ArchesTemplateEngine):
#         def document_replace(self, ):
#             pass
